"""
Content Format Adapters for LangGraph Knowledge Graph System

Provides adapter pattern for different content formats:
- PDF files
- Word documents (DOC, DOCX)
- PowerPoint presentations (PPT, PPTX)
- Text files (TXT, MD)
- HTML content
- JSON/XML structured data

Each adapter implements a common interface for content extraction and processing.
"""

import os
import logging
from abc import ABC, abstractmethod
from typing import Dict, Any, List, Optional, Union
from pathlib import Path
from datetime import datetime
import hashlib

# Content processing libraries
try:
    from langchain_community.document_loaders import PyPDFLoader, UnstructuredWordDocumentLoader, UnstructuredPowerPointLoader
    from langchain_community.document_loaders import TextLoader, UnstructuredHTMLLoader, JSONLoader
    CONTENT_LIBRARIES_AVAILABLE = True
except ImportError:
    CONTENT_LIBRARIES_AVAILABLE = False
    logging.warning("Content processing libraries not available. Install with: pip install unstructured python-docx python-pptx")

try:
    import pypdf
    import docx
    import pptx
    NATIVE_LIBRARIES_AVAILABLE = True
except ImportError:
    NATIVE_LIBRARIES_AVAILABLE = False
    logging.warning("Native content libraries not available. Install with: pip install pypdf python-docx python-pptx")

logger = logging.getLogger(__name__)

# ===============================
# ABSTRACT CONTENT ADAPTER
# ===============================

class ContentAdapter(ABC):
    """Abstract base class for content format adapters."""
    
    def __init__(self, file_path: str):
        self.file_path = Path(file_path)
        self.content_hash = self._generate_content_hash()
        self.metadata = self._extract_metadata()
    
    @abstractmethod
    def extract_text(self) -> str:
        """Extract plain text from the content."""
        pass
    
    @abstractmethod
    def extract_structured_content(self) -> Dict[str, Any]:
        """Extract structured content (headings, sections, etc.)."""
        pass
    
    @abstractmethod
    def get_content_type(self) -> str:
        """Get the content type identifier."""
        pass
    
    def _generate_content_hash(self) -> str:
        """Generate content hash for caching and versioning."""
        try:
            with open(self.file_path, 'rb') as f:
                content = f.read()
                return hashlib.md5(content).hexdigest()
        except Exception as e:
            logger.warning(f"Could not generate content hash: {e}")
            return f"hash_{datetime.now().isoformat()}"
    
    def _extract_metadata(self) -> Dict[str, Any]:
        """Extract basic file metadata."""
        try:
            stat = self.file_path.stat()
            return {
                "file_name": self.file_path.name,
                "file_size": stat.st_size,
                "created_time": datetime.fromtimestamp(stat.st_ctime).isoformat(),
                "modified_time": datetime.fromtimestamp(stat.st_mtime).isoformat(),
                "content_hash": self.content_hash
            }
        except Exception as e:
            logger.warning(f"Could not extract metadata: {e}")
            return {"file_name": self.file_path.name, "content_hash": self.content_hash}
    
    def get_processing_info(self) -> Dict[str, Any]:
        """Get comprehensive processing information."""
        return {
            "adapter_type": self.get_content_type(),
            "file_path": str(self.file_path),
            "metadata": self.metadata,
            "extraction_method": "langchain" if CONTENT_LIBRARIES_AVAILABLE else "native",
            "processed_at": datetime.now().isoformat()
        }

# ===============================
# PDF CONTENT ADAPTER
# ===============================

class PDFContentAdapter(ContentAdapter):
    """Adapter for PDF file processing."""
    
    def extract_text(self) -> str:
        """Extract text from PDF using LangChain or native library."""
        if CONTENT_LIBRARIES_AVAILABLE:
            try:
                loader = PyPDFLoader(str(self.file_path))
                pages = loader.load()
                return "\n".join([page.page_content for page in pages])
            except Exception as e:
                logger.error(f"LangChain PDF extraction failed: {e}")
        
        if NATIVE_LIBRARIES_AVAILABLE:
            try:
                text_content = []
                with open(self.file_path, 'rb') as file:
                    pdf_reader = pypdf.PdfReader(file)
                    for page in pdf_reader.pages:
                        text_content.append(page.extract_text())
                return "\n".join(text_content)
            except Exception as e:
                logger.error(f"Native PDF extraction failed: {e}")
        
        raise RuntimeError("No PDF extraction method available")
    
    def extract_structured_content(self) -> Dict[str, Any]:
        """Extract structured content from PDF."""
        text = self.extract_text()
        
        # Basic structure extraction (can be enhanced with AI)
        lines = text.split('\n')
        sections = []
        current_section = {"title": "Introduction", "content": []}
        
        for line in lines:
            line = line.strip()
            if not line:
                continue
            
            # Simple heading detection (can be enhanced)
            if line.isupper() and len(line) < 100:
                if current_section["content"]:
                    sections.append(current_section)
                current_section = {"title": line, "content": []}
            else:
                current_section["content"].append(line)
        
        if current_section["content"]:
            sections.append(current_section)
        
        return {
            "sections": sections,
            "total_pages": len(text.split('\n')) // 50,  # Rough estimate
            "word_count": len(text.split()),
            "structure_type": "document"
        }
    
    def get_content_type(self) -> str:
        return "pdf"

# ===============================
# WORD DOCUMENT ADAPTER
# ===============================

class WordContentAdapter(ContentAdapter):
    """Adapter for Word document processing."""
    
    def extract_text(self) -> str:
        """Extract text from Word document."""
        if CONTENT_LIBRARIES_AVAILABLE:
            try:
                loader = UnstructuredWordDocumentLoader(str(self.file_path))
                documents = loader.load()
                return "\n".join([doc.page_content for doc in documents])
            except Exception as e:
                logger.error(f"LangChain Word extraction failed: {e}")
        
        if NATIVE_LIBRARIES_AVAILABLE:
            try:
                doc = docx.Document(self.file_path)
                text_content = []
                for paragraph in doc.paragraphs:
                    text_content.append(paragraph.text)
                return "\n".join(text_content)
            except Exception as e:
                logger.error(f"Native Word extraction failed: {e}")
        
        raise RuntimeError("No Word document extraction method available")
    
    def extract_structured_content(self) -> Dict[str, Any]:
        """Extract structured content from Word document."""
        if NATIVE_LIBRARIES_AVAILABLE:
            try:
                doc = docx.Document(self.file_path)
                sections = []
                current_section = {"title": "Introduction", "content": [], "level": 0}
                
                for paragraph in doc.paragraphs:
                    text = paragraph.text.strip()
                    if not text:
                        continue
                    
                    # Detect headings based on style
                    if paragraph.style.name.startswith('Heading'):
                        if current_section["content"]:
                            sections.append(current_section)
                        level = int(paragraph.style.name.replace('Heading ', '')) if paragraph.style.name != 'Heading' else 1
                        current_section = {"title": text, "content": [], "level": level}
                    else:
                        current_section["content"].append(text)
                
                if current_section["content"]:
                    sections.append(current_section)
                
                return {
                    "sections": sections,
                    "total_paragraphs": len(doc.paragraphs),
                    "word_count": len(self.extract_text().split()),
                    "structure_type": "document"
                }
            except Exception as e:
                logger.warning(f"Structured Word extraction failed: {e}")
        
        # Fallback to basic text extraction
        text = self.extract_text()
        return {
            "sections": [{"title": "Content", "content": text.split('\n'), "level": 0}],
            "word_count": len(text.split()),
            "structure_type": "document"
        }
    
    def get_content_type(self) -> str:
        return "word"

# ===============================
# POWERPOINT ADAPTER
# ===============================

class PowerPointContentAdapter(ContentAdapter):
    """Adapter for PowerPoint presentation processing."""
    
    def extract_text(self) -> str:
        """Extract text from PowerPoint presentation."""
        if CONTENT_LIBRARIES_AVAILABLE:
            try:
                loader = UnstructuredPowerPointLoader(str(self.file_path))
                documents = loader.load()
                return "\n".join([doc.page_content for doc in documents])
            except Exception as e:
                logger.error(f"LangChain PowerPoint extraction failed: {e}")
        
        if NATIVE_LIBRARIES_AVAILABLE:
            try:
                prs = pptx.Presentation(self.file_path)
                text_content = []
                
                for slide in prs.slides:
                    slide_text = []
                    for shape in slide.shapes:
                        if hasattr(shape, "text"):
                            slide_text.append(shape.text)
                    if slide_text:
                        text_content.append(f"Slide {len(text_content) + 1}: {' '.join(slide_text)}")
                
                return "\n".join(text_content)
            except Exception as e:
                logger.error(f"Native PowerPoint extraction failed: {e}")
        
        raise RuntimeError("No PowerPoint extraction method available")
    
    def extract_structured_content(self) -> Dict[str, Any]:
        """Extract structured content from PowerPoint."""
        if NATIVE_LIBRARIES_AVAILABLE:
            try:
                prs = pptx.Presentation(self.file_path)
                slides = []
                
                for i, slide in enumerate(prs.slides):
                    slide_content = {
                        "slide_number": i + 1,
                        "title": "",
                        "content": [],
                        "notes": ""
                    }
                    
                    for shape in slide.shapes:
                        if hasattr(shape, "text"):
                            if shape.name.startswith("Title"):
                                slide_content["title"] = shape.text
                            else:
                                slide_content["content"].append(shape.text)
                    
                    slides.append(slide_content)
                
                return {
                    "slides": slides,
                    "total_slides": len(slides),
                    "word_count": len(self.extract_text().split()),
                    "structure_type": "presentation"
                }
            except Exception as e:
                logger.warning(f"Structured PowerPoint extraction failed: {e}")
        
        # Fallback to basic text extraction
        text = self.extract_text()
        return {
            "slides": [{"slide_number": 1, "title": "Content", "content": text.split('\n'), "notes": ""}],
            "word_count": len(text.split()),
            "structure_type": "presentation"
        }
    
    def get_content_type(self) -> str:
        return "powerpoint"

# ===============================
# TEXT CONTENT ADAPTER
# ===============================

class TextContentAdapter(ContentAdapter):
    """Adapter for text file processing."""
    
    def extract_text(self) -> str:
        """Extract text from text file."""
        if CONTENT_LIBRARIES_AVAILABLE:
            try:
                loader = TextLoader(str(self.file_path))
                documents = loader.load()
                return "\n".join([doc.page_content for doc in documents])
            except Exception as e:
                logger.error(f"LangChain text extraction failed: {e}")
        
        # Fallback to direct file reading
        try:
            with open(self.file_path, 'r', encoding='utf-8') as f:
                return f.read()
        except UnicodeDecodeError:
            try:
                with open(self.file_path, 'r', encoding='latin-1') as f:
                    return f.read()
            except Exception as e:
                logger.error(f"Text file reading failed: {e}")
                raise
    
    def extract_structured_content(self) -> Dict[str, Any]:
        """Extract structured content from text file."""
        text = self.extract_text()
        lines = text.split('\n')
        
        # Detect markdown-style headings
        sections = []
        current_section = {"title": "Content", "content": [], "level": 0}
        
        for line in lines:
            line = line.strip()
            if not line:
                continue
            
            # Markdown heading detection
            if line.startswith('#'):
                if current_section["content"]:
                    sections.append(current_section)
                level = len(line) - len(line.lstrip('#'))
                title = line.lstrip('#').strip()
                current_section = {"title": title, "content": [], "level": level}
            else:
                current_section["content"].append(line)
        
        if current_section["content"]:
            sections.append(current_section)
        
        return {
            "sections": sections,
            "total_lines": len(lines),
            "word_count": len(text.split()),
            "structure_type": "text"
        }
    
    def get_content_type(self) -> str:
        return "text"

# ===============================
# HTML CONTENT ADAPTER
# ===============================

class HTMLContentAdapter(ContentAdapter):
    """Adapter for HTML content processing."""
    
    def extract_text(self) -> str:
        """Extract text from HTML content."""
        if CONTENT_LIBRARIES_AVAILABLE:
            try:
                loader = UnstructuredHTMLLoader(str(self.file_path))
                documents = loader.load()
                return "\n".join([doc.page_content for doc in documents])
            except Exception as e:
                logger.error(f"LangChain HTML extraction failed: {e}")
        
        # Fallback to basic HTML parsing
        try:
            with open(self.file_path, 'r', encoding='utf-8') as f:
                html_content = f.read()
            
            # Simple HTML tag removal (can be enhanced with BeautifulSoup)
            import re
            text = re.sub(r'<[^>]+>', '', html_content)
            text = re.sub(r'\s+', ' ', text)
            return text.strip()
        except Exception as e:
            logger.error(f"HTML extraction failed: {e}")
            raise
    
    def extract_structured_content(self) -> Dict[str, Any]:
        """Extract structured content from HTML."""
        try:
            with open(self.file_path, 'r', encoding='utf-8') as f:
                html_content = f.read()
            
            # Basic HTML structure extraction
            import re
            
            # Extract headings
            headings = re.findall(r'<h[1-6][^>]*>(.*?)</h[1-6]>', html_content, re.IGNORECASE)
            
            # Extract paragraphs
            paragraphs = re.findall(r'<p[^>]*>(.*?)</p>', html_content, re.IGNORECASE)
            
            sections = []
            for i, heading in enumerate(headings):
                sections.append({
                    "title": heading.strip(),
                    "content": [paragraphs[i]] if i < len(paragraphs) else [],
                    "level": 1
                })
            
            return {
                "sections": sections,
                "headings": headings,
                "paragraphs": paragraphs,
                "word_count": len(self.extract_text().split()),
                "structure_type": "html"
            }
        except Exception as e:
            logger.warning(f"Structured HTML extraction failed: {e}")
            return {
                "sections": [{"title": "Content", "content": [self.extract_text()], "level": 0}],
                "word_count": len(self.extract_text().split()),
                "structure_type": "html"
            }
    
    def get_content_type(self) -> str:
        return "html"

# ===============================
# CONTENT ADAPTER FACTORY
# ===============================

class ContentAdapterFactory:
    """Factory for creating content adapters based on file type."""
    
    _adapters = {
        '.pdf': PDFContentAdapter,
        '.doc': WordContentAdapter,
        '.docx': WordContentAdapter,
        '.ppt': PowerPointContentAdapter,
        '.pptx': PowerPointContentAdapter,
        '.txt': TextContentAdapter,
        '.md': TextContentAdapter,
        '.html': HTMLContentAdapter,
        '.htm': HTMLContentAdapter,
    }
    
    @classmethod
    def create_adapter(cls, file_path: str) -> ContentAdapter:
        """Create appropriate adapter for the given file."""
        file_path = Path(file_path)
        extension = file_path.suffix.lower()
        
        if extension not in cls._adapters:
            raise ValueError(f"Unsupported file type: {extension}")
        
        adapter_class = cls._adapters[extension]
        return adapter_class(str(file_path))
    
    @classmethod
    def get_supported_formats(cls) -> List[str]:
        """Get list of supported file formats."""
        return list(cls._adapters.keys())
    
    @classmethod
    def is_supported(cls, file_path: str) -> bool:
        """Check if file format is supported."""
        extension = Path(file_path).suffix.lower()
        return extension in cls._adapters

# ===============================
# UNIFIED CONTENT PROCESSOR
# ===============================

class UnifiedContentProcessor:
    """Unified interface for processing any supported content format."""
    
    def __init__(self):
        self.factory = ContentAdapterFactory()
    
    def process_content(self, file_path: str) -> Dict[str, Any]:
        """Process content file and return unified structure."""
        try:
            # Create appropriate adapter
            adapter = self.factory.create_adapter(file_path)
            
            # Extract content
            text_content = adapter.extract_text()
            structured_content = adapter.extract_structured_content()
            processing_info = adapter.get_processing_info()
            
            return {
                "status": "success",
                "file_path": file_path,
                "content_type": adapter.get_content_type(),
                "text_content": text_content,
                "structured_content": structured_content,
                "processing_info": processing_info,
                "word_count": len(text_content.split()),
                "processed_at": datetime.now().isoformat()
            }
            
        except Exception as e:
            logger.error(f"Content processing failed for {file_path}: {e}")
            return {
                "status": "error",
                "file_path": file_path,
                "error": str(e),
                "processed_at": datetime.now().isoformat()
            }
    
    def get_supported_formats(self) -> List[str]:
        """Get list of supported file formats."""
        return self.factory.get_supported_formats()
    
    def is_format_supported(self, file_path: str) -> bool:
        """Check if file format is supported."""
        return self.factory.is_supported(file_path)

# ===============================
# USAGE EXAMPLES
# ===============================

def example_usage():
    """Example usage of content adapters."""
    processor = UnifiedContentProcessor()
    
    # Process different file types
    file_paths = [
        "document.pdf",
        "presentation.pptx", 
        "notes.docx",
        "readme.md"
    ]
    
    for file_path in file_paths:
        if processor.is_format_supported(file_path):
            result = processor.process_content(file_path)
            print(f"Processed {file_path}: {result['status']}")
        else:
            print(f"Unsupported format: {file_path}")

if __name__ == "__main__":
    example_usage() 